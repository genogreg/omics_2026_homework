from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "extra_credit_results"
FIG = OUT / "figures"
TAB = OUT / "tables"
PPTX = ROOT / "homeworks/final_extra_credit_presentation.pptx"


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)
    return box


def add_text(slide, lines, x=0.6, y=1.0, w=5.7, h=5.7, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(7)
    return box


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if w:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def metric_map():
    df = pd.read_csv(TAB / "promoter_summary.tsv", sep="\t")
    return {r.metric: r.value for r in df.itertuples(index=False)}


def make():
    metrics = metric_map()
    b = pd.read_csv(TAB / "boundary_parameter_summary.tsv", sep="\t")
    regions = pd.read_csv(TAB / "selected_regions.tsv", sep="\t")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = blank(prs)
    add_title(s, "Дополнительные задания по omics-анализу")
    add_text(s, [
        "Сделаны именно самостоятельные задания из конца материалов.",
        "Основной фокус: интеграция Hi-C, ChIP-seq, WGBS и RNAseq.",
        "Результаты: границы доменов, CTCF, промоторы, метилирование, активные гены.",
    ], x=0.8, y=1.4, w=11.5, h=3.5, size=22)

    s = blank(prs)
    add_title(s, "1. Параметры границ доменов")
    best = b.sort_values("center_enrichment", ascending=False).iloc[0]
    add_text(s, [
        f"Лучший вариант: {best.config}",
        f"Границ: {int(best.boundaries)}",
        f"С CTCF рядом +/-10 kb: {best.ctcf_near_fraction:.1%}",
        f"CTCF enrichment в центре: {best.center_enrichment:.2f}x",
        "Вывод: настройки меняют число границ, но CTCF остается обогащен около настоящих границ.",
    ], x=0.6, y=1.0, w=5.1, h=5.7)
    add_image(s, FIG / "boundary_parameters.png", 5.9, 1.05, w=6.8)

    s = blank(prs)
    add_title(s, "1. CTCF профиль около границ")
    add_text(s, [
        "Зеленые/цветные линии: CTCF signal около найденных границ.",
        "Черная пунктирная линия: случайный контроль.",
        "В центре границ CTCF выше контроля.",
        "Это подтверждает связь CTCF с частью Hi-C boundaries.",
    ], x=0.6, y=1.0, w=4.6, h=5.8)
    add_image(s, FIG / "ctcf_profiles_boundaries.png", 5.2, 1.0, w=7.5)

    s = blank(prs)
    add_title(s, "2. Промоторы: H3K27Ac, метилирование, RNAseq")
    add_text(s, [
        "Промотор: 1000 bp перед TSS.",
        "TSS приближен по CDS-аннотации.",
        f"Промоторов на chr1: {int(metrics['promoters_chr1'])}",
        f"С H3K27Ac peak: {int(metrics['promoters_with_h3k27ac_peak'])}",
        f"Низкое метилирование: {int(metrics['low_methylation_promoters'])}",
        f"H3K27Ac + low methylation + high RNAseq: {int(metrics['h3k27ac_low_methylation_active'])}",
        "Вывод: признаки активности согласуются.",
    ], x=0.6, y=1.0, w=5.3, h=6.0)
    add_image(s, FIG / "rnaseq_vs_methylation_h3k27ac.png", 6.1, 1.0, w=6.4)

    s = blank(prs)
    add_title(s, "4. Активные и неактивные промоторы")
    add_text(s, [
        f"Активные: RNAseq >= Q75, n={int(metrics['active_n'])}",
        f"Неактивные: RNAseq <= Q25, n={int(metrics['inactive_n'])}",
        f"H3K27Ac median: {metrics['active_h3k27ac_median']:.2f} vs {metrics['inactive_h3k27ac_median']:.2f}",
        f"Methylation median: {metrics['active_methylation_median']:.2f} vs {metrics['inactive_methylation_median']:.2f}",
        f"A-like compartment: {metrics['active_A_like_fraction']:.1%} vs {metrics['inactive_A_like_fraction']:.1%}",
        "Вывод: активные промоторы имеют больше H3K27Ac и ниже метилирование.",
    ], x=0.6, y=1.0, w=5.4, h=6.0)
    add_image(s, FIG / "active_inactive_promoters.png", 6.0, 1.0, w=6.8)

    s = blank(prs)
    add_title(s, "5. CTCF около активных генов")
    add_text(s, [
        f"Nearest CTCF median active: {metrics['active_ctcf_distance_median']:.0f} bp",
        f"Nearest CTCF median inactive: {metrics['inactive_ctcf_distance_median']:.0f} bp",
        "Разница небольшая.",
        "CTCF связан с архитектурой генома, но сам по себе не является прямой меткой активности.",
        "Главный вывод: связь активности с архитектурой умеренная.",
    ], x=0.7, y=1.15, w=11.8, h=5.0, size=20)

    s = blank(prs)
    add_title(s, "3. Регионы вместо IGV")
    add_text(s, [
        "IGV на сервере без GUI не использовался.",
        "Решение: pyGenomeTracks как терминальный аналог.",
        "Выбран пример активного гена HMGN2.",
        "Видно: RNAseq, CAGE, H3K27Ac, низкое метилирование, E1.",
    ], x=0.55, y=1.0, w=4.2, h=5.8)
    add_image(s, FIG / "genome_tracks_HMGN2.png", 4.85, 0.95, w=8.0)

    s = blank(prs)
    add_title(s, "3. Примеры активных регионов")
    lines = []
    for row in regions.itertuples(index=False):
        lines.append(f"{row.gene_name}: RNAseq {row.rnaseq:.1f}, methylation {row.methylation:.3f}, {row.compartment}")
    lines.append("Все три примера имеют H3K27Ac и низкое метилирование.")
    add_text(s, lines, x=0.75, y=1.1, w=11.8, h=5.5, size=20)

    s = blank(prs)
    add_title(s, "Вопросы WGBS")
    add_text(s, [
        "Low coverage CpG удаляются: при малом числе ридов beta-value нестабильна.",
        "Beta-value: доля метилированных ридов от 0 до 1.",
        "M-value: log2 отношение methylated / unmethylated.",
        f"High CpG promoters methylation: {metrics['high_cpg_promoter_methylation_median']:.3f}",
        f"Other promoters methylation: {metrics['other_promoter_methylation_median']:.3f}",
        "H3K27Ac связан с активными низкометилированными регуляторными участками.",
    ], x=0.75, y=1.0, w=11.8, h=6.0, size=19)

    s = blank(prs)
    add_title(s, "Итог")
    add_text(s, [
        "1. CTCF обогащен около Hi-C границ.",
        "2. H3K27Ac, низкое метилирование и RNAseq согласуются в активных промоторах.",
        "3. Активные промоторы чаще A-like и имеют ниже DNA methylation.",
        "4. CTCF-distance отличается слабо, значит архитектура связана с активностью не напрямую.",
        "5. Для серверной визуализации использован pyGenomeTracks вместо IGV.",
    ], x=0.75, y=1.0, w=11.8, h=5.8, size=21)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT

    PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX)
    print(PPTX)


if __name__ == "__main__":
    make()

